"""Render structured document content into real office files.

Consumed by environment materialization only — the simulation never sees
rendered bytes. Requires the ``artifacts`` extra: openpyxl for workbooks,
python-docx for documents, python-pptx for decks, reportlab for PDFs. No
external converter is involved, so every format works wherever the
package installs.
"""

from pathlib import Path

from pydantic import BaseModel, ConfigDict

from workbench.core.artifacts import (
    FormattedDocument,
    Formula,
    HeadingBlock,
    ListBlock,
    ParagraphBlock,
    SlideDeck,
    SpreadsheetContent,
    TableBlock,
    parse_formatted,
    parse_slides,
    parse_spreadsheet,
)


class RenderOutcome(BaseModel):
    model_config = ConfigDict(frozen=True, extra="forbid")

    path: Path
    # Human-readable reason when the requested form could not be produced
    # and a fallback was written instead. Never silently swallowed.
    skipped: str | None = None


def render_document(content_format: str, content: str, target: Path) -> RenderOutcome:
    """Write ``content`` to ``target`` in its real file form and return
    where it actually landed."""

    target.parent.mkdir(parents=True, exist_ok=True)
    match content_format:
        case "markdown":
            target.write_text(content, encoding="utf-8")
            return RenderOutcome(path=target)
        case "spreadsheet":
            try:
                parsed = parse_spreadsheet(content)
            except ValueError:
                return _raw_fallback(content_format, content, target)
            _render_xlsx(parsed, target)
            return RenderOutcome(path=target)
        case "formatted":
            try:
                document = parse_formatted(content)
            except ValueError:
                return _raw_fallback(content_format, content, target)
            if target.suffix.lower() == ".pdf":
                return _render_pdf(document, target)
            _render_docx(document, target)
            return RenderOutcome(path=target)
        case "slides":
            try:
                deck = parse_slides(content)
            except ValueError:
                return _raw_fallback(content_format, content, target)
            _render_pptx(deck, target)
            return RenderOutcome(path=target)
        case _:
            raise ValueError(f"unknown content format {content_format!r}")


def _raw_fallback(content_format: str, content: str, target: Path) -> RenderOutcome:
    """Authors sometimes declare a structured format but write prose; the
    record still deserves a file, so the raw text lands beside the intended
    name with the skip recorded."""

    fallback = target.with_suffix(".txt")
    fallback.write_text(content, encoding="utf-8")
    return RenderOutcome(
        path=fallback,
        skipped=(
            f"{target.name}: {content_format} content did not parse; "
            f"wrote raw text as {fallback.name}"
        ),
    )


def _render_xlsx(content: SpreadsheetContent, target: Path) -> None:
    from openpyxl import Workbook

    workbook = Workbook()
    workbook.remove(workbook.active)
    for sheet in content.sheets:
        worksheet = workbook.create_sheet(title=sheet.name)
        worksheet.append(list(sheet.columns))
        for row in sheet.rows:
            # openpyxl writes a leading-"=" string as a real formula, which
            # is what a workpaper needs: the file recalculates rather than
            # carrying a frozen number.
            worksheet.append(
                [cell.expression if isinstance(cell, Formula) else cell for cell in row]
            )
    workbook.save(target)


def _render_pptx(deck: SlideDeck, target: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    presentation = Presentation()
    title_and_body = presentation.slide_layouts[1]
    title_only = presentation.slide_layouts[5]
    for slide in deck.slides:
        layout = title_and_body if slide.bullets else title_only
        rendered = presentation.slides.add_slide(layout)
        rendered.shapes.title.text = slide.title
        if slide.bullets:
            frame = rendered.placeholders[1].text_frame
            frame.text = slide.bullets[0]
            for bullet in slide.bullets[1:]:
                paragraph = frame.add_paragraph()
                paragraph.text = bullet
        if slide.table is not None:
            columns = len(slide.table.columns)
            rows = len(slide.table.rows) + 1
            shape = rendered.shapes.add_table(
                rows,
                columns,
                Inches(0.5),
                Inches(4.0),
                Inches(9.0),
                Inches(0.4 * rows),
            )
            table = shape.table
            for index, heading in enumerate(slide.table.columns):
                table.cell(0, index).text = heading
            for row_index, row in enumerate(slide.table.rows, start=1):
                for column_index, value in enumerate(row):
                    cell = table.cell(row_index, column_index)
                    cell.text = str(value)
                    cell.text_frame.paragraphs[0].font.size = Pt(12)
        if slide.notes:
            rendered.notes_slide.notes_text_frame.text = slide.notes
    presentation.save(target)


def _render_docx(document: FormattedDocument, target: Path) -> None:
    from docx import Document

    output = Document()
    for block in document.blocks:
        match block:
            case HeadingBlock():
                output.add_heading(block.text, level=block.level)
            case ParagraphBlock():
                output.add_paragraph(block.text)
            case ListBlock():
                style = "List Number" if block.ordered else "List Bullet"
                for item in block.items:
                    output.add_paragraph(item, style=style)
            case TableBlock():
                table = output.add_table(
                    rows=len(block.rows) + 1, cols=len(block.columns)
                )
                for column, name in enumerate(block.columns):
                    table.rows[0].cells[column].text = name
                for index, row in enumerate(block.rows, start=1):
                    for column, value in enumerate(row):
                        table.rows[index].cells[column].text = value
    output.save(str(target))


def _render_pdf(document: FormattedDocument, target: Path) -> RenderOutcome:
    """Write a real PDF, with no external converter in the picture.

    This used to shell out to LibreOffice and, when `soffice` was absent,
    quietly leave a .docx where a .pdf was asked for — which is most of
    the time: not on this host, not in CI, not in the task container. A
    firm that never emits a PDF is not a firm.
    """

    from reportlab.lib import colors
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import (
        ListFlowable,
        ListItem,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )

    styles = getSampleStyleSheet()
    flow: list = []
    for block in document.blocks:
        match block:
            case HeadingBlock():
                level = min(max(block.level, 1), 4)
                flow.append(Paragraph(block.text, styles[f"Heading{level}"]))
            case ParagraphBlock():
                flow.append(Paragraph(block.text, styles["BodyText"]))
            case ListBlock():
                flow.append(
                    ListFlowable(
                        [
                            ListItem(Paragraph(item, styles["BodyText"]))
                            for item in block.items
                        ],
                        bulletType="1" if block.ordered else "bullet",
                    )
                )
            case TableBlock():
                data = [list(block.columns)] + [list(row) for row in block.rows]
                table = Table(data)
                table.setStyle(
                    TableStyle(
                        [
                            ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
                            ("BACKGROUND", (0, 0), (-1, 0), colors.whitesmoke),
                            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                            ("FONTSIZE", (0, 0), (-1, -1), 8),
                        ]
                    )
                )
                flow.append(table)
        flow.append(Spacer(1, 8))
    SimpleDocTemplate(str(target), pagesize=LETTER).build(flow)
    return RenderOutcome(path=target)
